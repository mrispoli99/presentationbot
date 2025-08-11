# app.py
import io
import os
import re
import zipfile
import tempfile
import shutil
import subprocess
from typing import List, Tuple, Dict, Optional

import streamlit as st
import pandas as pd

# ===============================
# Helpers (define BEFORE UI)
# ===============================

def clean_text(s):
    if s is None:
        return ""
    return re.sub(r"\s+", " ", str(s)).strip()

def first_row_as_header(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    first = df.iloc[0].astype(str).str.strip()
    nonempty_ratio = (first != "").mean()
    if nonempty_ratio >= 0.6 and first.nunique() >= max(2, int(df.shape[1] * 0.5)):
        df = df.copy()
        df.columns = [c if c else f"col_{i+1}" for i, c in enumerate(first)]
        df = df.iloc[1:].reset_index(drop=True)
    else:
        df.columns = [f"col_{i+1}" for i in range(df.shape[1])]
    return df

def sanitize_sheet_name(name: str) -> str:
    name = re.sub(r"[\\/*?:\[\]]", "_", name)[:31]
    return name or "Sheet"

def add_preview(name: str, df: pd.DataFrame):
    st.markdown(f"**{name}**")
    st.dataframe(df.head(30), use_container_width=True)

# ===============================
# Gemini helpers
# ===============================

def gemini_key():
    return os.getenv("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY", None)

def gemini_available() -> bool:
    try:
        return bool(gemini_key())
    except Exception:
        return False

def extract_text_from_gemini_response(resp) -> str:
    text = ""
    if hasattr(resp, "text") and isinstance(resp.text, str):
        text = (resp.text or "").strip()
    if not text and getattr(resp, "candidates", None):
        for c in resp.candidates:
            content = getattr(c, "content", None)
            parts = getattr(content, "parts", None) if content else None
            part_text = ""
            if parts:
                for p in parts:
                    if hasattr(p, "text") and p.text:
                        part_text += p.text
            if part_text:
                text = part_text.strip()
                break
    return text

def parse_csv_from_text(text: str) -> Optional[pd.DataFrame]:
    if not text:
        return None
    if "NO_TABLE" in text.upper():
        return None
    m = re.search(r"```(?:csv)?\s*(.*?)\s*```", text, flags=re.DOTALL | re.IGNORECASE)
    csv_text = m.group(1).strip() if m else text.strip()
    try:
        try:
            df = pd.read_csv(io.StringIO(csv_text))
        except Exception:
            df = pd.read_csv(io.StringIO(csv_text), sep=";")
        return df if not df.empty else None
    except Exception:
        return None

def gemini_extract_csv_from_image(image_bytes: bytes, model_name: str, debug_prefix: str = "") -> Tuple[Optional[pd.DataFrame], str]:
    key = gemini_key()
    if not key:
        return None, "No GEMINI_API_KEY"
    try:
        import google.generativeai as genai
    except ImportError:
        return None, "`google-generativeai` not installed"

    try:
        genai.configure(api_key=key)
        model = genai.GenerativeModel(model_name)
        prompt = (
            "You are a chart/table digitizer. "
            "If the image contains a data table or a simple bar/line/scatter chart, "
            "return ONLY CSV with headers on the first row. "
            "If you cannot reliably produce a table, reply exactly: NO_TABLE"
        )
        resp = model.generate_content(
            [prompt, {"mime_type": "image/png", "data": image_bytes}],
            request_options={"timeout": 90},
        )
        text = extract_text_from_gemini_response(resp)
        df = parse_csv_from_text(text)
        preview = text[:300].replace("\n", " ") if text else "<<empty>>"
        if df is None:
            return None, f"{debug_prefix} no CSV parsed; raw preview: {preview}"
        return df, f"{debug_prefix} CSV parsed OK; rows={len(df)}, cols={len(df.columns)}"
    except Exception as e:
        return None, f"Gemini exception: {type(e).__name__}: {e}"

# ===============================
# EMF/WMF → PNG conversion
# ===============================

def convert_emf_wmf_to_png_bytes(raw_bytes: bytes, ext: str) -> Tuple[Optional[bytes], str]:
    """
    Convert EMF/WMF bytes to PNG bytes using Inkscape (preferred) or ImageMagick.
    Returns (png_bytes or None, debug_message).
    """
    ext = ext.lower()
    if ext not in (".emf", ".wmf"):
        return None, f"skip: not EMF/WMF ({ext})"

    inkscape = shutil.which("inkscape") or shutil.which("inkscape.com")
    magick = shutil.which("magick")

    if not inkscape and not magick:
        return None, "no converter found (install Inkscape or ImageMagick)"

    src = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
    dst = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    try:
        src.write(raw_bytes)
        src.flush()
        src.close()
        dst.close()

        if inkscape:
            cmd = [inkscape, src.name, "--export-type=png", f"--export-filename={dst.name}"]
            try:
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
                with open(dst.name, "rb") as f:
                    return f.read(), "converted via Inkscape"
            except Exception as e:
                if not magick:
                    return None, f"Inkscape failed: {e}"

        if magick:
            cmd = [magick, src.name, dst.name]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
            with open(dst.name, "rb") as f:
                return f.read(), "converted via ImageMagick"

        return None, "conversion failed"
    finally:
        try: os.remove(src.name)
        except Exception: pass
        try: os.remove(dst.name)
        except Exception: pass

# ===============================
# PPTX: tables via python-pptx
# ===============================

def extract_pptx_tables(file_bytes: bytes, filename: str) -> List[Tuple[str, pd.DataFrame]]:
    results = []
    try:
        from pptx import Presentation
    except ImportError:
        st.info("`python-pptx` not installed; skipping PPT/PPTX table extraction.")
        return results

    prs = Presentation(io.BytesIO(file_bytes))
    for si, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            # breadth-first over groups
            queue = [sh]
            while queue:
                node = queue.pop(0)
                if getattr(node, "has_table", False):
                    tbl = node.table
                    rows = []
                    for r in tbl.rows:
                        rows.append([clean_text(c.text) for c in r.cells])
                    df = first_row_as_header(pd.DataFrame(rows))
                    results.append((f"{filename} - Slide {si} - Table", df))
                if hasattr(node, "shapes"):
                    queue.extend(list(node.shapes))
    return results

# ===============================
# PPTX: chart data (XML cache + embedded workbooks)
# ===============================

def _parse_chart_xml(xml_bytes: bytes) -> Optional[pd.DataFrame]:
    try:
        from lxml import etree
    except ImportError:
        return None
    ns = {
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    root = etree.fromstring(xml_bytes)
    data_dict: Dict[str, Dict[str, float]] = {}
    categories_set = set()

    for ser in root.findall(".//c:ser", namespaces=ns):
        ser_name = None
        tx_v = ser.find(".//c:tx//c:v", namespaces=ns)
        if tx_v is not None and tx_v.text:
            ser_name = clean_text(tx_v.text)
        if not ser_name:
            ser_name = "Series"

        cats = []
        cat_pts = ser.findall(".//c:cat//c:strCache//c:pt", namespaces=ns)
        if not cat_pts:
            cat_pts = ser.findall(".//c:cat//c:numCache//c:pt", namespaces=ns)
        for pt in cat_pts:
            v = pt.find("./c:v", namespaces=ns)
            cats.append(clean_text(v.text if v is not None else ""))

        vals = []
        val_pts = ser.findall(".//c:val//c:numCache//c:pt", namespaces=ns)
        for pt in val_pts:
            v = pt.find("./c:v", namespaces=ns)
            try:
                vals.append(float((v.text if v is not None else "nan")))
            except:
                vals.append(None)

        n = max(len(cats), len(vals))
        cats += [""] * (n - len(cats))
        vals += [None] * (n - len(vals))

        if ser_name not in data_dict:
            data_dict[ser_name] = {}
        for c, v in zip(cats, vals):
            categories_set.add(c)
            data_dict[ser_name][c] = v

    if not data_dict:
        return None

    cats_sorted = list(categories_set)
    first_ser = next(iter(data_dict))
    if data_dict[first_ser]:
        ordered = [c for c in data_dict[first_ser].keys() if c in categories_set]
        for c in cats_sorted:
            if c not in ordered:
                ordered.append(c)
        cats_sorted = ordered

    df = pd.DataFrame({"category": cats_sorted})
    for ser_name, ser_map in data_dict.items():
        df[ser_name] = [ser_map.get(c) for c in cats_sorted]
    return df

def extract_pptx_charts(file_bytes: bytes, filename: str) -> List[Tuple[str, pd.DataFrame]]:
    out = []
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
    except Exception:
        return out

    # cached chart XMLs
    for name in zf.namelist():
        if name.startswith("ppt/charts/") and name.endswith(".xml"):
            try:
                xml_bytes = zf.read(name)
                df = _parse_chart_xml(xml_bytes)
                if isinstance(df, pd.DataFrame) and not df.empty:
                    out.append((f"{filename} - {name.split('/')[-1]}", df))
            except Exception:
                continue

    # embedded workbooks
    for name in zf.namelist():
        if name.startswith("ppt/embeddings/") and (name.endswith(".xlsx") or name.endswith(".xls")):
            try:
                wb_bytes = zf.read(name)
                bio = io.BytesIO(wb_bytes)
                if name.endswith(".xlsx"):
                    xls = pd.ExcelFile(bio, engine="openpyxl")
                else:
                    try:
                        xls = pd.ExcelFile(bio, engine="xlrd")
                    except Exception:
                        continue
                for sheet in xls.sheet_names:
                    try:
                        df = xls.parse(sheet)
                        if not df.empty:
                            out.append((f"{filename} - EmbeddedWorkbook - {sheet}", df))
                    except Exception:
                        continue
            except Exception:
                continue
    return out

# ===============================
# PDF tables: Camelot -> pdfplumber
# ===============================

def extract_pdf_tables(file_bytes: bytes, filename: str) -> List[Tuple[str, pd.DataFrame]]:
    results = []
    used_camelot = False
    try:
        import camelot
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            tables = camelot.read_pdf(tmp_path, pages="all", flavor="lattice")
            for i, t in enumerate(tables or []):
                df = first_row_as_header(t.df)
                results.append((f"{filename} - Table (lattice) #{i+1}", df))
            used_camelot = bool(tables and len(tables) > 0)
        except Exception:
            pass
        if not used_camelot:
            try:
                tables = camelot.read_pdf(tmp_path, pages="all", flavor="stream")
                for i, t in enumerate(tables or []):
                    df = first_row_as_header(t.df)
                    results.append((f"{filename} - Table (stream) #{i+1}", df))
                used_camelot = bool(tables and len(tables) > 0)
            except Exception:
                pass
        try:
            os.remove(tmp_path)
        except Exception:
            pass
    except ImportError:
        pass

    if results:
        return results

    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            idx = 0
            for pi, page in enumerate(pdf.pages, start=1):
                try:
                    tables = page.extract_tables() or []
                except Exception:
                    tables = []
                for t in tables:
                    df = first_row_as_header(pd.DataFrame(t))
                    idx += 1
                    results.append((f"{filename} - Page {pi} - Table #{idx}", df))
    except ImportError:
        st.info("`camelot` and `pdfplumber` not installed; skipping PDF table extraction.")

    return results

# ===============================
# PDF: page raster → Gemini
# ===============================

def extract_pdf_charts_via_gemini(file_bytes: bytes, filename: str, model_name: str, dpi_scale: float = 2.5) -> List[Tuple[str, pd.DataFrame]]:
    results = []
    if not gemini_available():
        return results
    try:
        import fitz  # PyMuPDF
    except ImportError:
        st.warning("PyMuPDF (`fitz`) is not installed, so pages could not be rasterized for Gemini.")
        return results

    pages_attempted = 0
    pages_with_csv = 0
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for pi, page in enumerate(doc, start=1):
            mat = fitz.Matrix(dpi_scale, dpi_scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            pages_attempted += 1
            df, debug = gemini_extract_csv_from_image(img_bytes, model_name=model_name, debug_prefix=f"{filename} p{pi}:")
            st.write(debug)
            if isinstance(df, pd.DataFrame) and not df.empty:
                pages_with_csv += 1
                results.append((f"{filename} - Page {pi} - GeminiDigitized", df))

    st.write(f"Gemini scanned {pages_attempted} page(s) from **{filename}** → tables extracted: {pages_with_csv}.")
    return results

# ===============================
# PPTX: picture shapes → Gemini (PNG/JPG + EMF/WMF)
# ===============================

def extract_pptx_images_via_gemini(file_bytes: bytes, filename: str, model_name: str) -> List[Tuple[str, pd.DataFrame]]:
    """
    Finds images used on each slide and sends PNG/JPG (and EMF/WMF via conversion) to Gemini for CSV.
    """
    results = []
    if not gemini_available():
        return results

    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
    except Exception:
        return results

    slide_xml_names = sorted(
        [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")],
        key=lambda p: int(re.findall(r"slide(\d+)\.xml", p)[0])
    )

    import xml.etree.ElementTree as ET
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    for slide_path in slide_xml_names:
        slide_idx = int(re.findall(r"slide(\d+)\.xml", slide_path)[0])

        # relationships
        rels_path = slide_path.replace("slides/", "slides/_rels/") + ".rels"
        rels = {}
        if rels_path in zf.namelist():
            rels_xml = ET.fromstring(zf.read(rels_path))
            for rel in rels_xml.findall(".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
                rels[rel.attrib.get("Id")] = rel.attrib.get("Target")

        # referenced images
        slide_xml = ET.fromstring(zf.read(slide_path))
        blips = slide_xml.findall(".//a:blip", ns)

        img_targets = []
        for b in blips:
            rid = b.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            if not rid or rid not in rels:
                continue
            target = rels[rid]
            if target.startswith("../media/"):
                target = "ppt/media/" + target.split("../media/")[1]
            elif not target.startswith("ppt/"):
                target = "ppt/slides/" + target
            img_targets.append(target)

        for i, target in enumerate(img_targets, start=1):
            if target not in zf.namelist():
                st.write(f"{filename} s{slide_idx} img{i}: target missing {target}")
                continue
            ext = os.path.splitext(target.lower())[1]
            raw = zf.read(target)

            if ext in {".png", ".jpg", ".jpeg"}:
                png_bytes = raw  # raster
                method = "native raster"
            elif ext in {".emf", ".wmf"}:
                png_bytes, msg = convert_emf_wmf_to_png_bytes(raw, ext)
                if not png_bytes:
                    st.write(f"{filename} s{slide_idx} img{i}: {msg}")
                    continue
                method = msg
            else:
                st.write(f"{filename} s{slide_idx} img{i}: skipped unsupported {ext}")
                continue

            df, debug = gemini_extract_csv_from_image(
                png_bytes, model_name=model_name, debug_prefix=f"{filename} s{slide_idx} img{i} ({method}):"
            )
            st.write(debug)
            if isinstance(df, pd.DataFrame) and not df.empty:
                results.append((f"{filename} - Slide {slide_idx} - Image {i} - GeminiDigitized", df))

    return results

# ===============================
# Build outputs
# ===============================

def build_excel(dfs: List[Tuple[str, pd.DataFrame]]) -> bytes:
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        used_names = set()
        for name, df in dfs:
            sheet = sanitize_sheet_name(name[:31])
            base = sheet
            k = 1
            while sheet in used_names:
                sheet = sanitize_sheet_name(f"{base[:28]}_{k}")
                k += 1
            used_names.add(sheet)
            df.to_excel(writer, index=False, sheet_name=sheet)
    return output.getvalue()

def build_zip_csv(dfs: List[Tuple[str, pd.DataFrame]]) -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, mode="w", compression=zipfile.ZIP_DEFLATED) as z:
        used = set()
        for name, df in dfs:
            fname = re.sub(r"[^A-Za-z0-9_.-]+", "_", name)[:80]
            if not fname.lower().endswith(".csv"):
                fname = f"{fname}.csv"
            base = fname
            k = 1
            while fname in used:
                parts = base.rsplit(".", 1)
                fname = f"{parts[0]}_{k}.{parts[1]}"
                k += 1
            used.add(fname)
            z.writestr(fname, df.to_csv(index=False))
    return output.getvalue()

# ===============================
# Streamlit UI
# ===============================

st.set_page_config(page_title="Deck & PDF → Tables (CSV/Excel)", layout="wide")
st.title("📊 Slide/PDF Data Extractor → CSV/Excel")

with st.sidebar:
    st.header("Options")
    use_gemini_pdf = st.toggle(
        "Use Gemini for chart images in PDFs",
        value=True,
        help="Requires GEMINI_API_KEY. Converts page images to CSV (best-effort)."
    )
    use_gemini_ppt = st.toggle(
        "Use Gemini for images in PPT/PPTX (incl. EMF/WMF)",
        value=True,
        help="Sends slide images (PNG/JPG/EMF/WMF) to Gemini for CSV."
    )
    model_name = st.selectbox(
        "Gemini model",
        options=["gemini-2.5-flash", "gemini-2.0-flash-exp", "gemini-1.5-flash", "gemini-1.5-pro"],
        index=0
    )
    dpi_scale = st.slider("PDF page render scale (DPI factor)", 1.0, 3.0, 2.5, 0.5)
    st.divider()
    gem_ready = "✅" if gemini_available() else "❌"
    inkscape_ready = "✅" if (shutil.which("inkscape") or shutil.which("inkscape.com")) else "❌"
    magick_ready = "✅" if shutil.which("magick") else "❌"
    st.markdown(f"**Gemini detected:** {gem_ready}")
    st.markdown(f"Inkscape: {inkscape_ready} &nbsp;&nbsp; ImageMagick: {magick_ready}")
    if gem_ready == "❌":
        st.caption("Set GEMINI_API_KEY as env var or in `.streamlit/secrets.toml`.")
    if st.button("Test Gemini", use_container_width=True):
        if not gemini_available():
            st.error("No GEMINI_API_KEY found.")
        else:
            try:
                import google.generativeai as genai
                genai.configure(api_key=gemini_key())
                model = genai.GenerativeModel(model_name)
                r = model.generate_content("Reply with the single word: OK")
                txt = extract_text_from_gemini_response(r)
                st.success(f"Gemini response: {txt[:200] or '<<empty>>'}")
            except Exception as e:
                st.error(f"Gemini test failed: {type(e).__name__}: {e}")

st.write(
    "Upload **PPT / PPTX / PDF**. The app extracts native tables & chart data where available. "
    "For image-only slides/pages (including EMF/WMF), we rasterize and let Gemini estimate CSV."
)

uploads = st.file_uploader(
    "Upload one or more files",
    type=["ppt", "pptx", "pdf"],
    accept_multiple_files=True,
)

if st.button("Process files", type="primary", disabled=not uploads):
    all_results: List[Tuple[str, pd.DataFrame]] = []
    with st.status("Processing…", expanded=True) as status:
        for f in uploads:
            st.write(f"**{f.name}**")
            content = f.read()
            lower = f.name.lower()

            if lower.endswith((".ppt", ".pptx")):
                # Native tables
                tbls = extract_pptx_tables(content, f.name)
                st.write(f"- PowerPoint tables found: {len(tbls)}")
                all_results.extend(tbls)

                # Chart datasets (embedded/cached)
                ch = extract_pptx_charts(content, f.name)
                st.write(f"- PowerPoint chart datasets found: {len(ch)}")
                all_results.extend(ch)

                # Images → Gemini (PNG/JPG + EMF/WMF via converter)
                if use_gemini_ppt:
                    if not gemini_available():
                        st.warning("GEMINI_API_KEY not found; skipping PPT image → Gemini.")
                    else:
                        gem_tbls = extract_pptx_images_via_gemini(content, f.name, model_name=model_name)
                        st.write(f"- PPT images digitized by Gemini: {len(gem_tbls)}")
                        all_results.extend(gem_tbls)

            elif lower.endswith(".pdf"):
                # PDF vector tables
                pdf_tbls = extract_pdf_tables(content, f.name)
                st.write(f"- PDF vector tables found: {len(pdf_tbls)}")
                all_results.extend(pdf_tbls)

                # PDF pages → Gemini
                if use_gemini_pdf:
                    if not gemini_available():
                        st.warning("GEMINI_API_KEY not found; skipping PDF page images → Gemini.")
                    else:
                        gem_tbls = extract_pdf_charts_via_gemini(content, f.name, model_name=model_name, dpi_scale=dpi_scale)
                        all_results.extend(gem_tbls)

            else:
                st.warning(f"Unsupported file: {f.name}")

        status.update(label="Done", state="complete", expanded=False)

    if not all_results:
        st.warning(
            "No tables extracted. Checks: "
            "1) Click **Test Gemini**, "
            "2) `pip install -U PyMuPDF google-generativeai python-pptx`, "
            "3) For PDFs, increase **PDF page render scale**, "
            "4) Very low-res/hand-drawn charts may not parse."
        )
    else:
        st.subheader("Preview")
        for name, df in all_results[:20]:
            add_preview(name, df)

        st.subheader("Downloads")
        excel_bytes = build_excel(all_results)
        st.download_button(
            "⬇️ Download Excel (multi-sheet)",
            data=excel_bytes,
            file_name="extracted_data.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        zip_bytes = build_zip_csv(all_results)
        st.download_button(
            "⬇️ Download ZIP of CSVs",
            data=zip_bytes,
            file_name="extracted_csvs.zip",
            mime="application/zip",
        )

st.caption(
    "Notes: PPT Slide tables are read directly. Chart data is pulled from embedded workbooks/caches when present. "
    "Image-only slides/pages (PNG/JPG/EMF/WMF) are rasterized and sent to Gemini. "
    "Install Inkscape or ImageMagick for EMF/WMF support."
)
